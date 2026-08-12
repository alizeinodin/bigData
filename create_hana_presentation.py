#!/usr/bin/env python3
"""Generate Persian PowerPoint from SAP S/4HANA in-memory computing article."""

import os
from pathlib import Path

import arabic_reshaper
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from bidi.algorithm import get_display
from matplotlib import font_manager
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FONT_PATH = Path("/workspace/fonts/Vazir-Regular.ttf")
FONT_BOLD_PATH = Path("/workspace/fonts/Vazir-Bold.ttf")
CHARTS_DIR = Path("/workspace/presentation_charts")
FONT_NAME = "Vazir"

TITLE_COLOR = RGBColor(0, 51, 102)
ACCENT_COLOR = RGBColor(0, 102, 153)
BODY_COLOR = RGBColor(40, 40, 40)

STUDENT_NAME = "سید علی زین الدینی"
PROFESSOR_NAME = "دکتر صادق زاده"
COURSE_NAME = "پایگاه داده پیشرفته"


def fa(text: str) -> str:
    return get_display(arabic_reshaper.reshape(text))


def setup_matplotlib():
    font_manager.fontManager.addfont(str(FONT_PATH))
    font_manager.fontManager.addfont(str(FONT_BOLD_PATH))
    plt.rcParams["font.family"] = FONT_NAME
    plt.rcParams["axes.unicode_minus"] = False


def set_rtl(paragraph):
    try:
        from pptx.oxml.ns import qn

        pPr = paragraph._p.get_or_add_pPr()
        pPr.set(qn("a:rtl"), "1")
    except Exception:
        pass


def style_paragraph(paragraph, size=17, bold=False, color=BODY_COLOR, align=PP_ALIGN.RIGHT):
    paragraph.font.name = FONT_NAME
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.alignment = align
    set_rtl(paragraph)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(245, 248, 252)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.5))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "محاسبات درون‌حافظه‌ای با کارایی بالا"
    style_paragraph(tp, size=34, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.0))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "بررسی پژوهشی لایه پایگاه‌داده SAP S/4HANA"
    style_paragraph(sp, size=20, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)

    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(2.2))
    tf = info_box.text_frame
    tf.word_wrap = True
    info_lines = [
        f"ارائه‌دهنده: {STUDENT_NAME}",
        f"استاد درس: {PROFESSOR_NAME}",
        f"نام درس: {COURSE_NAME}",
        "",
        "منبع: Raikar, T. (2025) — American Journal of Technology",
        "DOI: 10.58425/ajt.v4i2.449",
    ]
    for i, line in enumerate(info_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        style_paragraph(p, size=18 if i < 3 else 14, bold=i < 3, color=TITLE_COLOR if i < 3 else BODY_COLOR, align=PP_ALIGN.CENTER)
        p.space_after = Pt(6)


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = TITLE_COLOR

    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = title
    style_paragraph(p, size=36, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


def add_slide_header(slide, title):
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    style_paragraph(tp, size=24, bold=True, color=TITLE_COLOR)

    line = slide.shapes.add_shape(1, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()


def add_content_slide(prs, title, bullets, sub_bullets=None, font_size=15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, title)

    body_box = slide.shapes.add_textbox(Inches(0.45), Inches(1.05), Inches(9.1), Inches(6.2))
    tf = body_box.text_frame
    tf.word_wrap = True
    sub_bullets = sub_bullets or {}

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        style_paragraph(p, size=font_size)
        p.space_after = Pt(6)

        if bullet in sub_bullets:
            for sub in sub_bullets[bullet]:
                sp = tf.add_paragraph()
                sp.text = "   ◦ " + sub
                style_paragraph(sp, size=font_size - 1, color=RGBColor(70, 70, 70))
                sp.level = 1


def add_content_with_image_slide(prs, title, bullets, image_path, image_left=5.0, font_size=14):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, title)

    text_width = image_left - 0.6
    body_box = slide.shapes.add_textbox(Inches(0.45), Inches(1.05), Inches(text_width), Inches(6.0))
    tf = body_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        style_paragraph(p, size=font_size)
        p.space_after = Pt(5)

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(image_left), Inches(1.2), width=Inches(10 - image_left - 0.3))


def add_image_slide(prs, title, image_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, title)

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.6), Inches(1.15), width=Inches(8.8))

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(9), Inches(0.5))
        cp = cap_box.text_frame.paragraphs[0]
        cp.text = caption
        style_paragraph(cp, size=12, color=RGBColor(90, 90, 90), align=PP_ALIGN.CENTER)


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, title)

    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.3), Inches(1.05), Inches(9.4), Inches(5.9)).table

    col_width = Inches(9.4 / cols)
    for c in range(cols):
        table.columns[c].width = col_width

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            style_paragraph(p, size=11, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                style_paragraph(p, size=10, align=PP_ALIGN.CENTER)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 245, 250)


def create_charts():
    CHARTS_DIR.mkdir(exist_ok=True)
    setup_matplotlib()
    colors = ["#003366", "#006699", "#3399cc", "#66b2ff", "#99ccff"]

    # 1) Architecture layers
    fig, ax = plt.subplots(figsize=(10, 6))
    layers = [
        (fa("کاربران و اپلیکیشن (Fiori / ABAP)"), "#003366"),
        (fa("Calculation Engine — SQL / SQL Script / MDX"), "#005588"),
        (fa("Optimizer + Execution Engine"), "#0077aa"),
        (fa("In-Memory Engines (Column / Graph / Text)"), "#0099cc"),
        (fa("Transaction Manager + Metadata Manager"), "#33aadd"),
        (fa("Persistence Layer (Log / Recovery / Disk)"), "#66ccee"),
    ]
    y = 0
    for label, color in layers:
        rect = mpatches.FancyBboxPatch((1.5, y), 7, 0.75, boxstyle="round,pad=0.02", facecolor=color, edgecolor="white")
        ax.add_patch(rect)
        ax.text(5, y + 0.375, label, ha="center", va="center", color="white", fontsize=11, fontweight="bold")
        y += 0.9
    ax.set_xlim(0, 10)
    ax.set_ylim(-0.2, y + 0.2)
    ax.axis("off")
    ax.set_title(fa("معماری لایه‌ای SAP S/4HANA"), fontsize=14, fontweight="bold", pad=12)
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "architecture_layers.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 2) Row vs Column data read
    fig, ax = plt.subplots(figsize=(9, 5.5))
    labels = [fa("Row Store\n(اسکن کامل)"), fa("Row Store\n(B-tree)"), fa("Column Store")]
    values = [7200, 360, 900]
    bars = ax.bar(labels, values, color=["#cc4444", "#dd8844", "#228844"], width=0.55)
    ax.set_ylabel(fa("حجم داده خوانده‌شده (MB)"), fontsize=11)
    ax.set_title(fa("مقایسه حجم I/O: Row Store در برابر Column Store"), fontsize=13, fontweight="bold")
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 120, f"{val:,} MB", ha="center", fontsize=10)
    ax.set_ylim(0, 8200)
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "row_vs_column.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 3) Column store benefits
    fig, ax = plt.subplots(figsize=(9, 5.5))
    metrics = [fa("کاهش اسکن"), fa("Dictionary"), fa("Vector"), fa("فشرده‌سازی"), fa("Aggregation"), fa("Parallel Scan")]
    speedups = [87, 3, 4.5, 6, 70, 10]
    bars = ax.barh(metrics, speedups, color=colors[:6])
    ax.set_xlabel(fa("ضریب بهبود (×) یا درصد"), fontsize=11)
    ax.set_title(fa("مزایای ذخیره‌سازی ستونی در SAP HANA"), fontsize=13, fontweight="bold")
    for bar, val in zip(bars, speedups):
        ax.text(bar.get_width() + 0.3, bar.get_y() + bar.get_height() / 2, f"{val}×", va="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "column_benefits.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 4) HTAP comparison
    fig, ax = plt.subplots(figsize=(9, 5.5))
    areas = [fa("گزارش"), fa("ETL"), fa("کوئری OLAP"), fa("Dashboard")]
    ecc = [12, 8, 150, 10]
    hana = [0.5, 0.1, 0.8, 0.2]
    x = range(len(areas))
    width = 0.35
    ax.bar([i - width / 2 for i in x], ecc, width, label=fa("ECC سنتی"), color="#cc4444")
    ax.bar([i + width / 2 for i in x], hana, width, label=fa("SAP HANA (HTAP)"), color="#228844")
    ax.set_ylabel(fa("زمان (ثانیه / ساعت — نسبی)"), fontsize=11)
    ax.set_xticks(list(x))
    ax.set_xticklabels(areas)
    ax.legend()
    ax.set_title(fa("مقایسه ECC و S/4HANA در بارهای OLTP+OLAP"), fontsize=13, fontweight="bold")
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "htap_comparison.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 5) Parallel processing
    fig, ax = plt.subplots(figsize=(9, 5))
    ops = [fa("Column Scan"), fa("Aggregation"), fa("Join"), fa("CPU Efficiency")]
    gains = [10, 8, 6, 75]
    ax.bar(ops, gains, color=["#003366", "#006699", "#3399cc", "#66b2ff"])
    ax.set_ylabel(fa("بهبود (× یا %)"), fontsize=11)
    ax.set_title(fa("مزایای پردازش موازی در HANA"), fontsize=13, fontweight="bold")
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "parallel_processing.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 6) Push-down flow
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.text(2, 3.5, fa("روش سنتی"), ha="center", fontsize=12, fontweight="bold")
    ax.annotate("", xy=(4.5, 3.5), xytext=(3, 3.5), arrowprops=dict(arrowstyle="->", lw=2))
    ax.text(5.5, 3.5, fa("App Server\n(ABAP Loop)"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#ffcccc"))
    ax.annotate("", xy=(8, 3.5), xytext=(6.5, 3.5), arrowprops=dict(arrowstyle="->", lw=2))
    ax.text(1.5, 3.5, fa("DB"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#ccccff"))
    ax.text(8.5, 3.5, fa("نتیجه"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#ccffcc"))

    ax.text(2, 1.5, fa("Push-Down"), ha="center", fontsize=12, fontweight="bold")
    ax.text(1.5, 1.5, fa("DB\n(HANA)"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#ccccff"))
    ax.annotate("", xy=(4, 1.5), xytext=(2.5, 1.5), arrowprops=dict(arrowstyle="->", lw=2, color="green"))
    ax.text(5.5, 1.5, fa("محاسبه در RAM\nJoin + Aggregate"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#ccffcc"))
    ax.annotate("", xy=(8, 1.5), xytext=(6.8, 1.5), arrowprops=dict(arrowstyle="->", lw=2, color="green"))
    ax.text(8.5, 1.5, fa("فقط\nنتیجه"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#ffffcc"))
    ax.text(5, 0.3, fa("کاهش ۸۰–۹۵٪ انتقال داده"), ha="center", fontsize=11, color="green")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 4.2)
    ax.axis("off")
    ax.set_title(fa("مقایسه جریان داده: سنتی در برابر Push-Down"), fontsize=13, fontweight="bold")
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "pushdown_flow.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 7) Query optimization
    fig, ax = plt.subplots(figsize=(9, 5.5))
    areas = [fa("انتقال"), fa("Loop"), fa("Aggregate"), fa("Join"), fa("Filter")]
    improvements = [87, 30, 40, 15, 12]
    ax.bar(areas, improvements, color=colors[:5])
    ax.set_ylabel(fa("بهبود (×)"), fontsize=11)
    ax.set_title(fa("مزایای بهینه‌سازی کوئری با Push-Down"), fontsize=13, fontweight="bold")
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "query_optimization.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 8) Dictionary + Vector diagram
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.text(2, 4, fa("جدول Orders"), ha="center", fontweight="bold")
    ax.text(2, 3.2, fa("Amount: 100,200,300,100,150"), ha="center", fontsize=10)
    ax.annotate("", xy=(5, 3.2), xytext=(3.5, 3.2), arrowprops=dict(arrowstyle="->", lw=2))
    ax.text(6.5, 3.5, fa("Dictionary\n100→0\n200→1\n300→2\n150→3"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#ddeeff"))
    ax.text(6.5, 1.5, fa("Vector\n[0,1,2,0,3]"), ha="center", va="center", bbox=dict(boxstyle="round", facecolor="#eeffdd"))
    ax.annotate("", xy=(6.5, 2.7), xytext=(6.5, 3.0), arrowprops=dict(arrowstyle="->", lw=2))
    ax.text(5, 0.5, fa("مصرف RAM کمتر + دسترسی سریع‌تر"), ha="center", fontsize=11, color="#006699")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 4.5)
    ax.axis("off")
    ax.set_title(fa("فشرده‌سازی Dictionary و Vector در Column Store"), fontsize=13, fontweight="bold")
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "dictionary_vector.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    # 9) Partitioning types
    fig, ax = plt.subplots(figsize=(10, 5))
    parts = [
        (fa("Hash Partitioning"), fa("توزیع یکنواخت بر اساس hash کلید\nمناسب: CUSTOMER_ID"), "#003366"),
        (fa("Range Partitioning"), fa("تقسیم بر اساس بازه (سال/ماه)\nPartition Pruning"), "#006699"),
        (fa("Round-Robin"), fa("توزیع چرخشی\nمناسب: staging / ETL"), "#3399cc"),
    ]
    for i, (title, desc, color) in enumerate(parts):
        x = 1.2 + i * 2.8
        rect = mpatches.FancyBboxPatch((x, 1.5), 2.4, 2.2, boxstyle="round,pad=0.05", facecolor=color, edgecolor="white")
        ax.add_patch(rect)
        ax.text(x + 1.2, 3.2, title, ha="center", va="center", color="white", fontsize=11, fontweight="bold")
        ax.text(x + 1.2, 2.2, desc, ha="center", va="center", color="white", fontsize=9)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis("off")
    ax.set_title(fa("انواع پارتیشن‌بندی در SAP HANA"), fontsize=13, fontweight="bold")
    fig.tight_layout()
    fig.savefig(CHARTS_DIR / "partitioning_types.png", dpi=150, bbox_inches="tight")
    plt.close(fig)

    return CHARTS_DIR


def build_presentation():
    charts = create_charts()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_content_slide(
        prs,
        "چکیده مقاله",
        [
            "هدف پژوهش: بررسی چگونگی دستیابی SAP S/4HANA به سرعت پردازش فوق‌العاده از طریق معماری درون‌حافظه‌ای و تحلیل لایه پایگاه‌داده HANA که منبع اصلی چالش‌های کارایی در سیستم‌های SAP است.",
            "روش کار: مطالعه مستندات SAP HANA و S/4HANA، اجرای کوئری‌های SQL، CDS Views و SQL Script روی سیستم واقعی، و تحلیل با ابزارهای ST05، ST12، SAT، DBACOCKPIT و PlanViz.",
            "یافته‌های کلیدی: ذخیره‌سازی ستونی، فشرده‌سازی Dictionary/Vector، پردازش Set-based، ایندکس‌گذاری پیشرفته و پارتیشن‌بندی، هرکدام به‌طور جداگانه و در کنار هم کارایی، مقیاس‌پذیری و تحلیل بلادرنگ را به‌شدت افزایش می‌دهند.",
            "نتیجه کلی: ترکیب طراحی درون‌حافظه‌ای با بهینه‌سازی‌های آگاهانه در لایه پایگاه‌داده، پایه و اساس پردازش بلادرنگ و تراکنش‌های زیر یک ثانیه در S/4HANA است.",
            "توصیه: توسعه‌دهندگان باید برنامه‌نویسی Set-based، ایندکس‌گذاری مناسب و پارتیشن‌بندی هوشمند را در اولویت قرار دهند و محاسبات را تا حد ممکن در لایه DB اجرا کنند.",
        ],
        font_size=14,
    )

    add_content_slide(
        prs,
        "کلمات کلیدی",
        [
            "SAP S/4HANA — نسل جدید ERP ساپ با موتور HANA",
            "In-Memory Computing — نگهداری و پردازش داده در RAM به‌جای دیسک",
            "Columnar Storage — ذخیره‌سازی ستونی برای کوئری‌های تحلیلی",
            "Real-time Analytics — تحلیل بلادرنگ بدون تأخیر ETL",
            "Set-based Processing — پردازش مجموعه‌ای به‌جای حلقه‌های تک‌سطری",
            "Push-down Processing — انتقال منطق محاسباتی به لایه پایگاه‌داده",
            "Parallel Processing & Table Partitioning — پردازش و توزیع موازی داده",
        ],
        font_size=15,
    )

    add_section_slide(prs, "۱. مقدمه")

    add_content_with_image_slide(
        prs,
        "چرا لایه پایگاه‌داده مهم است؟",
        [
            "کارایی SAP S/4HANA تحت تأثیر سه لایه وابسته است: پایگاه‌داده، لایه کاربرد (Application) و زیرساخت سخت‌افزاری. بهینه‌سازی هر لایه برای دستیابی به کارایی پایدار ضروری است.",
            "در پایگاه‌داده‌های سنتی مبتنی بر دیسک، هر کوئری باید داده را از storage فیزیکی بخواند که تأخیر I/O ایجاد می‌کند. این موضوع در سیستم‌هایی با هزاران کاربر همزمان به‌شدت محسوس می‌شود.",
            "S/4HANA کل جداول، ایندکس‌ها و نتایج میانی را در RAM نگه می‌دارد. نوشتن روی دیسک فقط برای persistence، backup و recovery انجام می‌شود.",
            "دسترسی به داده در میکروثانیه (به‌جای میلی‌ثانیه) تحول بنیادینی ایجاد می‌کند: تراکنش‌های زیر یک ثانیه، dashboard بلادرنگ و تحلیل روی داده عملیاتی.",
            "معماری HTAP (Hybrid OLTP+OLAP) امکان اجرای همزمان تراکنش و تحلیل را بدون نیاز به ETL شبانه فراهم می‌کند.",
        ],
        str(charts / "architecture_layers.png"),
        image_left=5.2,
        font_size=13,
    )

    add_image_slide(
        prs,
        "معماری SAP S/4HANA",
        str(charts / "architecture_layers.png"),
        "شکل ۱: لایه‌های اصلی معماری In-Memory SAP S/4HANA",
    )

    add_content_slide(
        prs,
        "اجزای کلیدی معماری (توضیحات)",
        [
            "Connection & Session Management: مدیریت اتصال کاربران و نشست‌ها — نقطه ورود تمام درخواست‌ها به سیستم.",
            "Authorization Manager: کنترل دسترسی — فقط کاربران مجاز می‌توانند داده را ببینند یا تغییر دهند.",
            "Calculation Engine: دریافت کوئری‌های SQL، SQL Script و MDX و انجام محاسبات.",
            "Optimizer & Plan Generator: یافتن سریع‌ترین plan اجرا با تحلیل cost-based.",
            "In-Memory Processing Engines: Column/Row Engine (جداول)، Graph Engine (شبکه/سلسله‌مراتب)، Text Engine (جستجوی متنی).",
            "Transaction Manager: تضمین ACID و consistency در بار همزمان.",
            "Persistence Layer: logging، recovery و ذخیره‌سازی پایدار برای جلوگیری از از دست رفتن داده.",
        ],
        font_size=14,
    )

    add_section_slide(prs, "۲. مرور ادبیات")

    add_content_slide(
        prs,
        "پژوهش‌های پیشین",
        [
            "Zhang Mei (2024): تأکید بر استفاده از مدل درون‌حافظه HANA برای تحلیل بلادرنگ، بهینه‌سازی مدل داده و کاهش افزونگی — نه فقط scale کردن سخت‌افزار.",
            "IRJMETS (2022/2025): تکنیک‌های partition pruning، push-down projection، data tiering و workload-aware tuning برای کاهش مصرف حافظه و runtime.",
            "Müller et al. (2015): بهینه‌سازی join با aggregate cache و object-awareness — بهبود چند-magnitude در کوئری‌های join-aggregate.",
            "Plattner (2014): پایه نظری column store، compression و pushdown — بنیان‌گذار مفاهیم In-Memory Data Management.",
            "Jain (n.d.): بهینه‌سازی indexing و partitioning برای workloadهای حجیم.",
        ],
        font_size=14,
    )

    add_content_slide(
        prs,
        "شکاف‌های پژوهشی و نوآوری این مقاله",
        [
            "ادبیات موجود بیشتر ویژگی‌های HANA را توضیح می‌دهد اما معیارهای runtime واقعی در S/4HANA را کمتر گزارش کرده است.",
            "تأثیر indexing و partitioning روی workloadهای مختلف (OLTP، OLAP، mixed) به‌طور کمی بررسی نشده.",
            "ابزارهای SAP مانند ST05، PlanViz و DBACOCKPIT در پژوهش‌ها کمتر استفاده شده‌اند.",
            "چارچوب یکپارچه‌ای برای تعامل query design، compression، indexing و partitioning وجود ندارد.",
            "این مقاله با رویکرد عملی (hands-on) و استفاده از ابزارهای trace، این شکاف‌ها را پر می‌کند و یک framework یکپارچه ارائه می‌دهد.",
        ],
        font_size=14,
    )

    add_section_slide(prs, "۳. روش‌شناسی")

    add_content_slide(
        prs,
        "روش پژوهش",
        [
            "مبانی نظری: مستندات SAP HANA، راهنمای توسعه S/4HANA و ادبیات HTAP و in-memory databases.",
            "آزمایش عملی: اجرای SQL queries، CDS Views و SQL Script procedures روی SAP HANA با datasetهای کنترل‌شده.",
            "مقایسه: منطق ABAP مبتنی بر LOOP AT ... ENDLOOP در برابر اجرای Set-based و push-down در DB.",
            "ابزارهای تحلیل: ST05 (SQL trace)، ST12 (ABAP+DB trace)، SAT (runtime analysis)، DBACOCKPIT (DB monitoring)، PlanViz (execution plan visualization).",
            "معیارها: execution plan، CPU usage، scan volume، memory usage، push-down efficiency و parallel processing behavior.",
            "هدف: درک اینکه چگونه storage structures، compression، indexing و partitioning روی runtime تأثیر می‌گذارند.",
        ],
        font_size=14,
    )

    add_section_slide(prs, "۴. اصول مدیریت داده درون‌حافظه‌ای")

    add_content_with_image_slide(
        prs,
        "ذخیره‌سازی ستونی (Column Store)",
        [
            "در Row Store داده‌ها سطر به سطر ذخیره می‌شوند؛ در Column Store هر ستون جداگانه در حافظه قرار می‌گیرد.",
            "مزیت اصلی: کوئری فقط ستون‌های مورد نیاز را می‌خواند — کاهش چشمگیر I/O و CPU.",
            "Dictionary Encoding: مقادیر یکتا (distinct) در جدول lookup نگهداری می‌شوند؛ مقادیر تکراری ذخیره نمی‌شوند.",
            "Vector Encoding: به‌جای مقدار خام، ValueID (اشاره‌گر به dictionary) ذخیره می‌شود.",
            "این دو تکنیک هم فشرده‌سازی و هم سرعت lookup را بهبود می‌دهند.",
            "برای کوئری‌های SELECT SUM، GROUP BY و فیلتر روی ستون‌های محدود، Column Store به‌شدت برتر است.",
        ],
        str(charts / "dictionary_vector.png"),
        image_left=5.0,
        font_size=12,
    )

    add_image_slide(
        prs,
        "فشرده‌سازی Dictionary و Vector",
        str(charts / "dictionary_vector.png"),
        "شکل ۲: نحوه فشرده‌سازی ستون Amount با Dictionary و Vector",
    )

    add_image_slide(
        prs,
        "مزایای ذخیره‌سازی ستونی",
        str(charts / "column_benefits.png"),
        "شکل ۳: مقایسه بهبود سرعت تکنیک‌های Column Store",
    )

    add_table_slide(
        prs,
        "جدول مزایای Column Store (از مقاله)",
        ["بهینه‌سازی", "بهبود سرعت"],
        [
            ["Column Store — داده اسکن‌شده", "۸۰–۹۵٪ کاهش"],
            ["Dictionary Encoding", "۳× سریع‌تر"],
            ["Vector Encoding", "۳–۶× سریع‌تر"],
            ["فشرده‌سازی", "۳–۱۰× کوچک‌تر"],
            ["SUM/Aggregation Pushdown", "۶۰–۸۰٪ سریع‌تر"],
            ["Parallel Column Scans", "۵–۱۶× سریع‌تر"],
        ],
    )

    add_content_with_image_slide(
        prs,
        "مقایسه Row Store و Column Store",
        [
            "سناریو: جدول ORDERS با ۱۰۰ میلیون سطر — کوئری SUM(Amount) WHERE Country='IN'",
            "Row Store (اسکن کامل): ≈ ۷.۲ GB خوانده می‌شود؛ در حالی که کوئری فقط به یک ستون نیاز دارد — ۹۰٪ داده بی‌فایده است.",
            "Row Store (B-tree index روی Country): ≈ ۳۶۰ MB + overhead دسترسی تصادفی و cache miss.",
            "Column Store: فقط Country (≈100 MB) + Amount (≈800 MB) = ≈ ۹۰۰ MB — داده فشرده و sequential.",
            "نتیجه: Column Store حدود ۸× سریع‌تر و بسیار CPU-friendlyتر است.",
            "دلیل: اسکن sequential ستونی + فشرده‌سازی + پردازش vectorized.",
        ],
        str(charts / "row_vs_column.png"),
        image_left=4.8,
        font_size=12,
    )

    add_image_slide(
        prs,
        "نمودار مقایسه I/O",
        str(charts / "row_vs_column.png"),
        "شکل ۴: حجم داده خوانده‌شده — Row Store در برابر Column Store",
    )

    add_content_with_image_slide(
        prs,
        "پردازش موازی (Parallel Processing)",
        [
            "SAP HANA از multi-core CPU و distributed nodes برای توزیع بار استفاده می‌کند.",
            "یک کوئری بزرگ به subtaskهای کوچک‌تر تقسیم و همزمان اجرا می‌شود — مثلاً هر region روی یک core.",
            "Parallel column scans: ۵–۱۶× سریع‌تر از اسکن تک‌رشته‌ای.",
            "Parallel aggregation: ۴–۱۲× سریع‌تر — SUM/COUNT روی partitionهای موازی.",
            "Parallel join: ۳–۱۰× سریع‌تر — hash/merge join موازی.",
            "CPU utilization: بهبود ۶۰–۹۰٪ — runtime از دقیقه به ثانیه.",
        ],
        str(charts / "parallel_processing.png"),
        image_left=5.0,
        font_size=12,
    )

    add_image_slide(
        prs,
        "HTAP: ترکیب OLTP و OLAP",
        str(charts / "htap_comparison.png"),
        "شکل ۵: مقایسه ECC سنتی و SAP HANA در بارهای ترکیبی",
    )

    add_table_slide(
        prs,
        "جدول HTAP: ECC در برابر S/4HANA",
        ["حوزه", "ECC سنتی", "SAP HANA", "بهبود"],
        [
            ["گزارش عملیاتی", "تأخیر ساعتی/روزانه", "۰–۱ ثانیه", "تا ۱۰۰۰×"],
            ["ETL/Batch", "شبانه/ساعتی", "حذف کامل", "۱۰۰٪"],
            ["کوئری تحلیلی", "۱۰–۳۰۰ ثانیه", "۰.۱–۱.۵ ثانیه", "۲۰–۲۰۰×"],
            ["Dashboard", "بعد از ETL", "بلادرنگ", "فوری"],
            ["پیچیدگی", "ECC+BW+ETL", "یک HANA", "۴۰–۶۰٪ کمتر"],
        ],
    )

    add_content_slide(
        prs,
        "HTAP / Translytical Processing (توضیحات)",
        [
            "در سیستم‌های قدیمی: OLTP (تراکنش) و OLAP (تحلیل) روی دو DB جدا — داده شبانه به BW منتقل می‌شد.",
            "نتیجه: dashboard همیشه دیروز را نشان می‌داد، نه وضعیت لحظه‌ای.",
            "S/4HANA هر دو را روی یک پلتفرم in-memory اجرا می‌کند — Translytical Processing.",
            "مثال: با ثبت سفارش فروش (OLTP)، مدیر بلافاصله گزارش سود و موجودی (OLAP) را می‌بیند.",
            "ETL حذف می‌شود → کاهش latency، ساده‌سازی landscape و هزینه عملیاتی.",
        ],
        font_size=14,
    )

    add_section_slide(prs, "۵. Push-Down Processing")

    add_image_slide(
        prs,
        "پردازش Push-Down",
        str(charts / "pushdown_flow.png"),
        "شکل ۶: جریان داده در روش سنتی در برابر Push-Down",
    )

    add_content_slide(
        prs,
        "پردازش Push-Down — توضیحات کامل",
        [
            "روش سنتی SAP: منطق کسب‌وکار در Application Server (ABAP) اجرا می‌شد. DB فقط store/retrieve می‌کرد.",
            "مشکل: SELECT * → انتقال میلیون‌ها سطر به App Server → LOOP AT → SELECT SINGLE — بار شبکه و CPU بالا.",
            "Push-Down: منطق به DB «فرو رانده» می‌شود — HANA join، aggregate و filter را در RAM انجام می‌دهد.",
            "فقط نتیجه نهایی (مثلاً یک سطر SUM) به App Server برمی‌گردد — کاهش ۸۰–۹۵٪ transfer.",
            "پیاده‌سازی: CDS Views، AMDP (ABAP Managed DB Procedures)، SQL Script.",
            "HANA به‌عنوان columnar + parallel engine، محاسبات را بسیار سریع‌تر از ABAP loop انجام می‌دهد.",
        ],
        font_size=14,
    )

    add_content_slide(
        prs,
        "طراحی کوئری (Query Design) — بهترین شیوه‌ها",
        [
            "به‌جای ABAP Loop از JOIN، CDS View و set-based logic استفاده کنید.",
            "CDS Views لایه‌بندی: Interface (I_) برای raw data، Composite (C_) برای join/enrichment، Consumption (Z_/Y_) برای Fiori/OData.",
            "فیلتر را زود اعمال کنید: WHERE قبل از HAVING — کاهش داده قبل از join.",
            "Window Functions: ROW_NUMBER، RANK، LEAD، LAG — محاسبات analytical بدون GROUP BY.",
            "Associations به‌جای join مستقیم — خوانایی و filter pushdown بهتر.",
            "Clean Core: namespace Z/Y، Extension Views، Adapt before Build — سازگاری با upgrade.",
        ],
        font_size=14,
    )

    add_image_slide(
        prs,
        "مزایای بهینه‌سازی کوئری",
        str(charts / "query_optimization.png"),
        "شکل ۷: بهبود سرعت با Push-Down در HANA",
    )

    add_table_slide(
        prs,
        "جدول بهینه‌سازی کوئری",
        ["حوزه", "ECC / ABAP", "HANA Push-Down", "بهبود"],
        [
            ["انتقال داده", "داده خام زیاد", "فقط نتیجه", "۸۰–۹۵٪"],
            ["Loop", "App Server", "Vectorized HANA", "۱۰–۵۰×"],
            ["Aggregate", "Loop ABAP", "Columnar RAM", "۲۰–۶۰×"],
            ["Join", "SELECT SINGLE", "Parallel join", "۵–۳۰×"],
            ["Filter", "Boolean loop", "ValueID pushdown", "۵–۲۰×"],
        ],
    )

    add_content_slide(
        prs,
        "ایندکس‌گذاری (Indexing) — انواع و کاربرد",
        [
            "Dictionary encoding در column store خود یک index سبک است — ValueID-based filtering.",
            "Inverted Index: نگاشت ValueID → Row IDs — جستجوی سریع مقادیر (مثلاً Artist = Lady Gaga → rows 2,6).",
            "Full-Text Index: جستجوی متنی — حذف stop words و inverted index برای terms.",
            "Multi-Column Index: dictionary ترکیبی برای چند ستون (Region + Product).",
            "Inverted Value: سریع‌ترین read/join — PK و unique constraints.",
            "Inverted Hash: ≈۳۰٪ کمتر RAM — PK با ستون‌های بلند.",
            "Inverted Individual: DML سریع‌تر — read کمی کندتر (۱۰–۲۰٪).",
        ],
        font_size=13,
    )

    add_content_slide(
        prs,
        "بهترین شیوه‌های Index",
        [
            "Composite key: ستون‌های selective (بیشترین تمایز) اول، سپس equality، سپس range/sort.",
            "Join columns: ترتیب و data type یکسان در هر دو طرف join — مثلاً (BUKRS, BELNR).",
            "از index روی ستون‌های low-selectivity (Boolean، status با ۲–۳ مقدار) خودداری کنید.",
            "اگر PK/unique index همه predicates را پوشش می‌دهد، secondary index اضافه نکنید.",
            "برای workload تحلیلی scan-heavy: CDS logic، partitioning یا preaggregation — نه index زیاد.",
            "Expensive Statement Trace → PlanViz → ST05: فقط indexهای واقعاً لازم بسازید.",
        ],
        font_size=14,
    )

    add_image_slide(
        prs,
        "انواع پارتیشن‌بندی",
        str(charts / "partitioning_types.png"),
        "شکل ۸: Hash، Range و Round-Robin Partitioning",
    )

    add_content_slide(
        prs,
        "پارتیشن‌بندی جداول — توضیحات",
        [
            "Hash: CUSTOMER_ID → hash → partition — توزیع یکنواخت، عبور از limit ۲B rows، parallel scan.",
            "Range: FISCAL_YEAR → P2019, P2020, ... — partition pruning: فقط partition مرتبط اسکن می‌شود.",
            "Round-Robin: row1→P1, row2→P2, ... — balance load، مناسب staging/ETL، نه key lookup.",
            "Range + Hash composite: برای hot partitions (داده اخیر بیشتر) — range اول، سپس hash.",
            "مزایا: تا ۹۵٪ کمتر data scanned، ۵–۲۵× parallel scan، ۳–۱۵× join/aggregate.",
        ],
        font_size=14,
    )

    add_table_slide(
        prs,
        "جدول مزایای پارتیشن‌بندی",
        ["حوزه", "ECC", "SAP HANA", "بهبود"],
        [
            ["Range Query", "اسکن کل جدول", "Partition pruning", "تا ۹۵٪"],
            ["Parallel", "تک‌رشته‌ای", "هر partition روی core", "۵–۲۵×"],
            ["Join/Aggregate", "کل جدول", "Partition-wise", "۳–۱۵×"],
        ],
    )

    add_section_slide(prs, "۶. نتیجه‌گیری")

    add_content_slide(
        prs,
        "نتیجه‌گیری",
        [
            "معماری in-memory S/4HANA با column store، compression، parallel processing و HTAP، پارادigm پردازش enterprise را متحول کرده است.",
            "Push-down، indexing و partitioning باید به‌صورت یکپارچه و متناسب با workload طراحی شوند — نه به‌صورت جداگانه.",
            "انتقال محاسبات از ABAP Loop به CDS/SQL Script: کاهش چشمگیر runtime، network load و app server CPU.",
            "ابزارهای trace (ST05، PlanViz، DBACOCKPIT) برای performance engineering عملی ضروری هستند.",
            "Clean-core CDS practices: سیستم upgrade-safe و cloud-ready باقی می‌ماند.",
            "پژوهش آینده: سناریوهای صنعتی، multi-tenant cloud و AI workloads.",
        ],
        font_size=14,
    )

    add_content_slide(
        prs,
        "توصیه‌های عملی برای توسعه‌دهندگان",
        [
            "Set-based programming و push-down را به ABAP loop ترجیح دهید.",
            "Indexing و partitioning را بر اساس الگوی query و data volume طراحی کنید.",
            "CDS Views لایه‌بندی‌شده (I_/C_/Z_) و Clean Core را رعایت کنید.",
            "Performance engineering را از ابتدای چرخه پیاده‌سازی بگنجانید — نه فقط در انتها.",
            "از HTAP برای حذف ETL، dashboard بلادرنگ و analytics روی داده عملیاتی استفاده کنید.",
            "Expensive Statement Trace و PlanViz را برای شناسایی bottleneckها به‌کار ببرید.",
        ],
        font_size=14,
    )

    add_content_slide(
        prs,
        "منابع و اختصارات",
        [
            "CDS: Core Data Services | RAM: Random Access Memory | CPU: Central Processing Unit",
            "OLTP: Online Transaction Processing | OLAP: Online Analytical Processing",
            "ETL: Extract, Transform, Load | HTAP: Hybrid Transactional/Analytical Processing",
            "PK: Primary Key | AMDP: ABAP Managed Database Procedures",
            "Raikar, T. (2025). High-Performance In-Memory Computing: A Research Study on SAP S/4 HANA Database Layer. American Journal of Technology, 4(2), 94–113.",
            "DOI: https://doi.org/10.58425/ajt.v4i2.449",
        ],
        font_size=13,
    )

    # Closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(245, 248, 252)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    p = box.text_frame.paragraphs[0]
    p.text = "سپاس از توجه شما"
    style_paragraph(p, size=36, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    sub = box.text_frame.add_paragraph()
    sub.text = f"{STUDENT_NAME}\nسوالات؟"
    style_paragraph(sub, size=20, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    output = "/workspace/SAP_S4HANA_InMemory_Presentation_FA.pptx"
    presentation = build_presentation()
    presentation.save(output)
    print(f"Saved: {output}")
    print(f"Slides: {len(presentation.slides)}")
